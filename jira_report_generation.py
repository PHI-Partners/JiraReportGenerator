from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from datetime import datetime
from models import JiraReport

def get_month_year():
    return datetime.today().strftime("%B %Y")

def generate_pptx(template_pptx, jira_data, output_filename="Jira_Report.pptx"):
    prs = Presentation(template_pptx)
    jira_report = JiraReport.from_json(jira_data)

    # --- First Slide ---
    first_slide = prs.slides[0]
    # Loop through shapes in the first slide
    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            shape.text = f"Fusion Support Report\n{get_month_year()}"
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = 1
                for run in paragraph.runs:
                    # Specify the style of the text (e.g., font size, color, bold, etc.)
                    run.font.size = Pt(41)  # Set font size (320,000 EMU = 41 pt)
                    run.font.bold = True  # Make the text bold
                    run.font.color.rgb = RGBColor(47, 97, 174)  # Set text color to blue
                    run.font.name = 'Calibri'  # Set font family

    # --- Second Slide ---
    second_slide = prs.slides[1]
    new_header_text = f"Fusion Support Update ({get_month_year()})"

    # Define target subtitle texts
    target_subtitles = {
        "left": "Tickets distribution per Business Unit (day of report)",
        "right": "Tickets status (day of the report)"
    }

    # Store positions of detected subtitles
    subtitle_positions = {}

    for shape in second_slide.shapes:
        if hasattr(shape, "text"):
            if "Fusion Support Update" in shape.text:
                shape.text = new_header_text  # Update the text
            # Detect left subtitle
            elif shape.text.strip() == target_subtitles["left"]:
                subtitle_positions["left"] = {
                    "top": shape.top,
                    "left": shape.left,
                    "height": shape.height,
                    "width": shape.width
                }
            # Detect right subtitle
            elif shape.text.strip() == target_subtitles["right"]:
                subtitle_positions["right"] = {
                    "top": shape.top,
                    "left": shape.left,
                    "height": shape.height,
                    "width": shape.width
                }

    # --- Left Section: Business Unit Image (keep existing) ---
    if "left" in subtitle_positions:
        left_shape = subtitle_positions["left"]
        # Add image below left subtitle (your existing image code)
        # ...

        # --- Right Section: Status Text ---
    if "right" in subtitle_positions:
        right_shape = subtitle_positions["right"]

        # Create formatted status text
        status_text = "Status\t\t\tCount\n"  # Header

        for status, count in jira_report.statuses.items():
            status_text += f"{status}\t\t\t{count}\n"

        # Add textbox below right subtitle
        textbox = second_slide.shapes.add_textbox(
            left=right_shape["left"],
            top=right_shape["top"] + right_shape["height"] + Inches(0.2),
            width=right_shape["width"],
            height=Inches(2)  # Adjust height as needed
        )

        # Format the text
        tf = textbox.text_frame
        tf.text = status_text

        # Style the text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            paragraph.font.name = 'Calibri'

            # Bold the header row
            if paragraph.text.startswith("Status"):
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(47, 97, 174)  # Your blue color

    '''
    # --- Add Content Below Subtitles ---
    # Left Section (Business Unit Distribution)
    if "left" in subtitle_positions:
        pos = subtitle_positions["left"]
        # Add image 0.2 inches below subtitle
        second_slide.shapes.add_picture(
            "gym.png",  # Replace with actual image path
            left=pos["left"],
            top=pos["top"] + pos["height"] + Inches(0.2),
            width=Inches(4.5),  # Adjust as needed
            height=Inches(3.5)
        )
    # Right Section (Ticket Status)
    if "right" in subtitle_positions:
        pos = subtitle_positions["right"]
        # Add image 0.2 inches below subtitle
        second_slide.shapes.add_picture(
            "gym.png",  # Replace with actual image path
            left=pos["left"],
            top=pos["top"] + pos["height"] + Inches(0.2),
            width=Inches(4.5),
            height=Inches(3.5)
        )
    '''

    '''
    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]  # Title & Subtitle
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Jira Ticket Report"
    subtitle.text = "Generated on: 2025-04-02"

    # --- Summary Slide ---
    summary_slide_layout = prs.slide_layouts[1]  # Title & Content
    slide = prs.slides.add_slide(summary_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Jira Ticket Summary"
    content.text = f"Total Tickets: {jira_data['total']}"
    '''
    # ------------------------ Save PowerPoint ------------------------
    prs.save(output_filename)
    print(f"PowerPoint report saved: {output_filename}")

